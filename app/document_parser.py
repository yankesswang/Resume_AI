import os
import re
import gc
import time
import torch
import logging
from pathlib import Path
from typing import Dict, Tuple, Optional
from contextlib import contextmanager
from PIL import Image

# Marker PDF 相關
from marker.converters.pdf import PdfConverter
from marker.models import create_model_dict
from marker.output import text_from_rendered

logger = logging.getLogger(__name__)


class DocumentParser:
    """
    Unified document parser for:
    - PDF   (Marker span / image based)
    - DOCX  (Heading + image based)
    - PPTX  (Slide based)

    All formats output ordered chunks into:
        output_dir/ordered_splits/*.txt
    """

    # ============================================================
    # Basic utilities
    # ============================================================

    @contextmanager
    def _timed(self, name: str):
        """
        Simple timing context manager for profiling phases.
        """
        start = time.perf_counter()
        yield
        print(f"[Timing] {name}: {time.perf_counter() - start:.2f}s")

    def _cleanup(self):
        """
        Aggressive memory cleanup (important for Marker / CUDA).
        """
        gc.collect()
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
            torch.cuda.synchronize()

    def _safe(self, s: str) -> str:
        """
        Make string filesystem-safe for filenames.
        """
        return "".join(c for c in s if c.isalnum() or c in "-_")

    def _prepare_dirs(self, output_dir: str) -> str:
        """
        Ensure ordered_splits directory exists.
        """
        splits_dir = os.path.join(output_dir, "ordered_splits")
        os.makedirs(splits_dir, exist_ok=True)
        return splits_dir

    def _write(self, path: str, text: str):
        """
        Write UTF-8 text file.
        """
        with open(path, "w", encoding="utf-8") as f:
            f.write(text.strip())

    # ============================================================
    # PDF parsing (Marker-based)
    # ============================================================

    def parse_pdf(
        self,
        pdf_path: str,
        output_dir: str,
        pdftext_workers: int = 1,
        force_ocr: bool = False,
        return_images: bool = True,
        save_images: bool = True,
    ) -> Tuple[str, str, Dict[str, Image.Image]]:
        """
        Parse PDF using Marker.
        - Marker generates <span id=...> and image markdown
        - We split text based on those markers

        Returns:
            Tuple of (text, md_path, images_dict)
        """

        pdf_path = os.path.abspath(pdf_path)
        pdf_stem = Path(pdf_path).stem
        os.makedirs(output_dir, exist_ok=True)
        splits_dir = self._prepare_dirs(output_dir)

        # -------- Phase 1: Marker extraction --------
        converter = None
        rendered = None
        images = None
        try:
            with self._timed("PDF Marker extraction"):
                converter = PdfConverter(
                    artifact_dict=create_model_dict(),
                    config={
                        "pdftext_workers": pdftext_workers,
                        "force_ocr": force_ocr,
                    },
                )
                rendered = converter(pdf_path)
                text, _, images = text_from_rendered(rendered)
        finally:
            if rendered is not None:
                del rendered
            if converter is not None:
                del converter
            self._cleanup()

        if not text:
            raise RuntimeError("Marker returned empty text")

# -------- Phase 2: ordered splitting (span id + page + image) --------
        pattern = (
            r'(<span id="([^"]+)"></span>)'
            r'|(!\[.*?\]\(([^)]*?_page_(\d+)_Figure_\d+\.jpeg)\))'
        )
        matches = list(re.finditer(pattern, text))

        counter = 0

        current_page = None
        page_span_index = {}

        def write_chunk(kind: str, name: str, content: str):
            nonlocal counter
            if not content.strip():
                return
            counter += 1
            fname = f"{counter:04d}_{kind}_{self._safe(name)}.txt"
            self._write(os.path.join(splits_dir, fname), content.strip())

        # ---------- header ----------
        if matches:
            header = text[:matches[0].start()]
            write_chunk("span", "header", header)

        # ---------- main loop ----------
        for i, m in enumerate(matches):
            start = m.end()
            end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
            between = text[start:end]

            # ===== span marker =====
            if m.group(2):
                span_id = m.group(2)

                # page marker: <span id="pageX">
                if span_id.startswith("page") and span_id[4:].isdigit():
                    current_page = int(span_id[4:])
                    page_span_index.setdefault(current_page, 0)
                    page_span_index[current_page] += 1
                    span_name = f"page{current_page}-{page_span_index[current_page]}"
                else:
                    # original semantic span id (abstract / sec:intro / etc.)
                    span_name = span_id

                write_chunk("span", span_name, between)

            # ===== image marker =====
            elif m.group(3):
                img_path = m.group(3)
                img_stem = Path(img_path).stem

                # image itself is an ordering node
                write_chunk("img", img_stem, f"![]({img_path})")

                # text after image → fallback to page span
                if current_page is not None:
                    page_span_index.setdefault(current_page, 0)
                    page_span_index[current_page] += 1
                    span_name = f"page{current_page}-{page_span_index[current_page]}"
                    write_chunk("span", span_name, between)

        # ---------- save images ----------
        images = images or {}
        if save_images:
            for name, img in images.items():
                img.save(os.path.join(output_dir, name))
        if not return_images:
            for img in images.values():
                try:
                    img.close()
                except Exception:
                    pass
            images = {}

        # ---------- save original markdown ----------
        md_path = os.path.join(output_dir, f"{pdf_stem}_original.md")
        self._write(md_path, text)

        self._cleanup()
        return text, md_path, images



    # ============================================================
    # Candidate splitting
    # ============================================================

    def split_candidates(self, markdown: str) -> list[str]:
        """
        Split a combined markdown (multiple resumes in one PDF)
        into individual candidate sections.

        Each candidate section starts with '# 基本資料' and ends with
        the privacy notice line.
        """
        # Split on '# 基本資料', '## 基本資料', '### 基本資料', etc.
        parts = re.split(r'(?=^#{1,6} 基本資料$)', markdown, flags=re.MULTILINE)
        # Filter out empty/whitespace-only parts (e.g. before the first match)
        candidates = [p.strip() for p in parts if p.strip()]
        return candidates

    # ============================================================
    # DOCX parsing (Heading + Image)
    # ============================================================

    def parse_docx(self, docx_path: str, output_dir: str) -> str:
        """
        Parse DOCX into ordered chunks.
        Chunk boundaries:
        - Heading paragraphs
        - Inline images (w:drawing)
        """

        import docx
        from docx.oxml.ns import qn

        doc = docx.Document(docx_path)
        splits_dir = self._prepare_dirs(output_dir)

        chunks = []           # [(span_id, content)]
        buffer = []           # current text buffer
        current_span = "start"
        image_index = 0

        def flush():
            """
            Flush buffered text into a chunk.
            """
            nonlocal buffer, current_span
            if buffer:
                chunks.append((current_span, "\n\n".join(buffer)))
                buffer = []

        for para in doc.paragraphs:
            text = para.text.strip()

            # -------- Heading handling --------
            if para.style.name.startswith("Heading") and text:
                flush()
                level = para.style.name.replace("Heading ", "")
                current_span = f"h{level}_{text}"
                buffer.append(f"{'#' * int(level)} {text}")
                continue

            # -------- Image handling --------
            # Images are stored inside w:r -> w:drawing
            for run in para._element.iter(qn("w:r")):
                drawing = run.find(qn("w:drawing"))
                if drawing is not None:
                    flush()

                    blip = drawing.find(
                        ".//a:blip",
                        namespaces={
                            "a": "http://schemas.openxmlformats.org/drawingml/2006/main"
                        },
                    )
                    if blip is None:
                        continue

                    rId = blip.get(qn("r:embed"))
                    part = doc.part.related_parts[rId]

                    image_index += 1
                    img_name = f"docx_image_{image_index:04d}.png"
                    img_path = os.path.join(output_dir, img_name)

                    with open(img_path, "wb") as f:
                        f.write(part.blob)

                    # Image itself is a chunk
                    chunks.append((f"img_{img_name}", f"![]({img_name})"))

            # -------- Normal paragraph text --------
            if text:
                buffer.append(text)

        flush()

        # -------- Write ordered_splits --------
        for i, (span_id, content) in enumerate(chunks, 1):
            kind = "img" if span_id.startswith("img_") else "span"
            fname = f"{i:04d}_{kind}_{self._safe(span_id)}.txt"
            self._write(os.path.join(splits_dir, fname), content)

        return "\n\n".join(c[1] for c in chunks)

    # ============================================================
    # PPTX parsing (Slide-based)
    # ============================================================

    def parse_pptx(self, pptx_path: str, output_dir: str) -> str:
        """
        Parse PPTX.
        Each slide becomes one ordered chunk, with images extracted.
        """

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(pptx_path)
        splits_dir = self._prepare_dirs(output_dir)
        os.makedirs(output_dir, exist_ok=True)

        chunks = []
        image_index = 0

        for idx, slide in enumerate(prs.slides, 1):
            buffer = [f"## Slide {idx}"]

            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, "text") and shape.text.strip():
                    buffer.append(shape.text.strip())

                # Extract images in-place to preserve order
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_index += 1
                        ext = image.ext or "png"
                        img_name = f"pptx_slide{idx}_image_{image_index:04d}.{ext}"
                        img_path = os.path.join(output_dir, img_name)

                        with open(img_path, "wb") as f:
                            f.write(image.blob)

                        buffer.append(f"![]({img_name})")
                    except Exception as e:
                        logger.warning(f"Failed to extract image from slide {idx}: {e}")

            chunks.append((f"slide_{idx}", "\n\n".join(buffer)))

        for i, (sid, content) in enumerate(chunks, 1):
            fname = f"{i:04d}_span_{sid}.txt"
            self._write(os.path.join(splits_dir, fname), content)

        return "\n\n---\n\n".join(c[1] for c in chunks)

    # ============================================================
    # Unified entry point
    # ============================================================

    def convert(self, path: str, output_dir: Optional[str] = None) -> str:
        """
        Auto-detect file type and parse accordingly.
        """
        ext = Path(path).suffix.lower()
        output_dir = output_dir or os.path.dirname(path)

        if ext == ".pdf":
            text, _, _ = self.parse_pdf(path, output_dir)
            return text
        elif ext in (".docx", ".doc"):
            return self.parse_docx(path, output_dir)
        elif ext in (".pptx", ".ppt"):
            return self.parse_pptx(path, output_dir)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def convert_to_markdown(self, file_path: str, file_extension: str) -> str:
        """
        Convert document to Markdown format (compatibility method).

        Args:
            file_path: Path to the document file
            file_extension: File extension (.pdf, .docx, etc.)

        Returns:
            Markdown formatted text
        """
        output_dir = os.path.dirname(file_path)

        if file_extension == ".pdf":
            text, _, _ = self.parse_pdf(file_path, output_dir)
            return text
        elif file_extension in [".docx", ".doc"]:
            return self.parse_docx(file_path, output_dir)
        elif file_extension in [".pptx", ".ppt"]:
            return self.parse_pptx(file_path, output_dir)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

    def cleanup(self):
        """
        Public cleanup method for releasing resources.
        """
        self._cleanup()



if __name__ == "__main__":
    parser = DocumentParser()

    file_path = "data/履歷預覽.pdf"
    output_directory = "output/many_people"

    # 確保輸出目錄存在
    os.makedirs(output_directory, exist_ok=True)

    # 轉換並切 ordered_splits
    text = parser.convert(file_path, output_directory)

    print("=== Conversion Done ===")
    print(f"Input file: {file_path}")
    print(f"Output dir: {output_directory}")
    print(f"Extracted text length: {len(text)} characters")

    # 分割不同候選人
    candidates = parser.split_candidates(text)
    print(f"\n=== Found {len(candidates)} candidates ===")
    for i, c in enumerate(candidates, 1):
        # 從第一個表格取姓名
        name_match = re.search(r'姓/名:\s*\|\s*(\S+)', c)
        name = name_match.group(1) if name_match else f"candidate_{i}"
        candidate_dir = os.path.join(output_directory, f"candidate_{i}_{name}")
        os.makedirs(candidate_dir, exist_ok=True)
        md_path = os.path.join(candidate_dir, f"{name}.md")
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(c)
        print(f"  Candidate {i}: {name} ({len(c)} chars) -> {md_path}")
